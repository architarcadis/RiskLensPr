"""Export module for RiskLens Pro

Provides functionality to export reports and visualizations to PDF and PowerPoint formats.
"""

import os
import tempfile
import base64
import pandas as pd
import numpy as np
import streamlit as st
import plotly.express as px
import plotly.graph_objects as go
from fpdf import FPDF, HTMLMixin
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import io

# Define Arcadis brand colors
ARCADIS_PRIMARY_COLOR = "#FF6900"  # Orange
ARCADIS_SECONDARY_COLOR = "#4D4D4F"  # Dark gray
ARCADIS_ACCENT_COLOR = "#0063A3"  # Blue
ARCADIS_LIGHT_BG = "#F7F7F7"  # Light gray background
ARCADIS_SUCCESS = "#28A745"  # Green
ARCADIS_WARNING = "#FFC107"  # Yellow
ARCADIS_DANGER = "#DC3545"  # Red


class RiskLensPDF(FPDF):
    """Custom PDF class for RiskLens Pro reports"""
    
    def __init__(self, orientation='P', unit='mm', format='A4'):
        super().__init__(orientation, unit, format)
        self.set_auto_page_break(auto=True, margin=15)
        # Use standard fonts instead of custom ones
        self.set_font('Arial', 'B', 12)
        self.title = "RiskLens Pro Report"
        self.project_name = ""
        self.PRIMARY_COLOR = ARCADIS_PRIMARY_COLOR
        self.SECONDARY_COLOR = ARCADIS_SECONDARY_COLOR
        self.ACCENT_COLOR = ARCADIS_ACCENT_COLOR
    
    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB"""
        h = hex_color.lstrip('#')
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
    
    def header(self):
        """Page header"""
        # Arcadis logo
        if os.path.exists('./assets/arcadis_logo.png'):
            self.image('./assets/arcadis_logo.png', x=10, y=8, w=60)
        
        # Set header font
        self.set_font('Arial', 'B', 10)
        self.set_text_color(80, 80, 80)
        
        # Add report title
        self.cell(0, 10, self.title, 0, 0, 'R')
        
        # Add project name if available
        if self.project_name:
            self.ln(5)
            self.set_font('Arial', '', 9)
            self.set_x(10)
            self.cell(0, 10, f"Project: {self.project_name}", 0, 0, 'R')
        
        # Add date
        self.ln(5)
        self.set_font('Arial', '', 9)
        self.set_x(10)
        self.cell(0, 10, f"Generated: {pd.Timestamp.now().strftime('%B %d, %Y')}", 0, 0, 'R')
        
        # Add line
        self.ln(15)
        rgb = self.hex_to_rgb(self.PRIMARY_COLOR)
        self.set_draw_color(rgb[0], rgb[1], rgb[2])
        self.set_line_width(0.5)
        self.line(10, 25, 200, 25)
        self.ln(10)
    
    def footer(self):
        """Page footer"""
        # Set footer position
        self.set_y(-15)
        
        # Set font
        self.set_font('Arial', '', 8)
        self.set_text_color(80, 80, 80)
        
        # Add page number
        self.cell(0, 10, f"Page {self.page_no()}", 0, 0, 'C')
        
        # Add company info
        self.set_y(-10)
        self.cell(0, 10, "RiskLens Pro by Arcadis", 0, 0, 'C')
    
    def chapter_title(self, title, level=1):
        """Add a chapter title"""
        rgb = self.hex_to_rgb(self.PRIMARY_COLOR)
        self.set_text_color(rgb[0], rgb[1], rgb[2])
        
        # Set font size based on heading level
        if level == 1:
            self.set_font('Arial', 'B', 16)
        elif level == 2:
            self.set_font('Arial', 'B', 14)
        else:
            self.set_font('Arial', 'B', 12)
        
        # Add title
        self.cell(0, 10, title, 0, 1)
        self.ln(4)
        
        # Reset text color
        self.set_text_color(0, 0, 0)
    
    def chapter_body(self, text):
        """Add chapter text with automatic line breaks"""
        self.set_font('Arial', '', 11)
        self.multi_cell(0, 5, text)
        self.ln()
    
    def add_metric(self, label, value, icon=None, color=None):
        """Add a metric box"""
        # Default color if not specified
        if color is None:
            color = self.PRIMARY_COLOR
        
        rgb = self.hex_to_rgb(color)
        
        # Set starting position
        start_x = self.get_x()
        start_y = self.get_y()
        
        # Draw metric box
        self.set_fill_color(245, 245, 245)
        self.set_draw_color(rgb[0], rgb[1], rgb[2])
        self.set_line_width(0.3)
        self.rect(start_x, start_y, 60, 25, 'DF')
        
        # Add label
        self.set_font('Arial', '', 9)
        self.set_text_color(100, 100, 100)
        self.set_xy(start_x + 2, start_y + 3)
        self.cell(56, 5, label, 0, 1, 'L')
        
        # Add value
        self.set_text_color(rgb[0], rgb[1], rgb[2])
        self.set_font('Arial', 'B', 14)
        self.set_xy(start_x + 2, start_y + 12)
        self.cell(56, 8, str(value), 0, 1, 'C')
        
        # Reset position outside the box
        self.set_xy(start_x + 65, start_y)
        
        # Reset text color
        self.set_text_color(0, 0, 0)
    
    def add_figure(self, fig, width=190, height=100):
        """Add a plotly figure to the PDF by converting it to an image"""
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                tmp_filename = tmp.name
                
            # Save figure as image
            fig.write_image(tmp_filename, width=width*5, height=height*5)
            
            # Add image to PDF
            self.image(tmp_filename, x=10, y=self.get_y(), w=width)
            
            # Remove temporary file
            os.unlink(tmp_filename)
            
            # Add some space after the figure
            self.ln(height + 10)
            
        except Exception as e:
            # Fallback to text description if image creation fails
            self.set_font('Arial', 'B', 10)
            self.set_text_color(100, 100, 100)
            self.multi_cell(0, 8, "[Chart visualization available in the app]")
            
            # Add a description if the figure has a title
            if hasattr(fig, 'layout') and hasattr(fig.layout, 'title') and fig.layout.title:
                title = fig.layout.title.text if hasattr(fig.layout.title, 'text') else str(fig.layout.title)
                self.set_font('Arial', '', 9)
                self.multi_cell(0, 5, f"Chart title: {title}")
            
            # Add some space after the figure note
            self.ln(5)
    
    def add_table(self, df, max_rows=20, title=None):
        """Add a table to the PDF"""
        # Check if table fits on current page, otherwise add a new page
        if self.get_y() + 40 + (min(len(df), max_rows) * 8) > self.h - 20:
            self.add_page()
        
        # Add title if provided
        if title:
            self.set_font('Arial', 'B', 12)
            rgb = self.hex_to_rgb(self.PRIMARY_COLOR)
            self.set_text_color(rgb[0], rgb[1], rgb[2])
            self.cell(0, 10, title, 0, 1, 'L')
            self.ln(2)
        
        # Reset text color
        self.set_text_color(0, 0, 0)
        
        # Table header
        columns = df.columns
        col_widths = [180 / len(columns)] * len(columns)  # Distribute width evenly
        
        # Adjust specific column widths if needed
        for i, col in enumerate(columns):
            if 'ID' in col or 'Id' in col:
                col_widths[i] = max(col_widths[i] * 0.7, 15)  # Narrower for IDs
            elif 'Name' in col or 'Title' in col:
                col_widths[i] = min(col_widths[i] * 1.5, 70)  # Wider for names/titles
        
        # Normalize widths to ensure they sum to 180
        total_width = sum(col_widths)
        col_widths = [w * 180 / total_width for w in col_widths]
        
        # Draw header
        self.set_font('Arial', 'B', 10)
        rgb = self.hex_to_rgb(self.SECONDARY_COLOR)
        self.set_fill_color(rgb[0], rgb[1], rgb[2])
        self.set_text_color(255, 255, 255)
        
        for i, col in enumerate(columns):
            self.cell(col_widths[i], 8, str(col), 1, 0, 'C', True)
        self.ln()
        
        # Draw rows
        self.set_font('Arial', '', 9)
        self.set_text_color(0, 0, 0)
        
        # Handle too many rows
        display_df = df.head(max_rows) if len(df) > max_rows else df
        fill = False
        
        for _, row in display_df.iterrows():
            # Alternate row colors
            if fill:
                self.set_fill_color(240, 240, 240)
            else:
                self.set_fill_color(255, 255, 255)
            
            for i, col in enumerate(columns):
                value = row[col]
                # Format specific column types
                if isinstance(value, (np.float64, float)):
                    formatted_value = f"{value:.2f}" if abs(value) < 100 else f"{value:.1f}"
                elif isinstance(value, (np.int64, int)):
                    formatted_value = f"{value:,}"
                else:
                    formatted_value = str(value)[:25]  # Truncate long text
                
                self.cell(col_widths[i], 8, formatted_value, 1, 0, 'L', fill)
            
            self.ln()
            fill = not fill
        
        # Add note if table was truncated
        if len(df) > max_rows:
            self.set_font('Arial', '', 8)
            self.set_text_color(100, 100, 100)
            self.cell(0, 6, f"Showing {max_rows} of {len(df)} rows", 0, 1, 'L')
        
        self.ln(5)


def create_pdf_report(project_data, risk_analysis=None, visualizations=None, model_results=None, risk_data=None, report_type="detailed"):
    """Create a PDF report with project risk analysis
    
    Args:
        project_data: DataFrame with project data
        risk_analysis: Dictionary with risk analysis results
        visualizations: Dictionary with plotly figures
        model_results: Dictionary of model results
        risk_data: DataFrame with risk data
        report_type: Type of report to generate (executive, detailed, etc.)
        
    Returns:
        bytes: PDF report as bytes object
    """
    # Create PDF document
    pdf = RiskLensPDF()
    pdf.title = f"RiskLens Pro {report_type.title()} Report"
    
    # Set project name if available
    if project_data is not None and 'ProjectName' in project_data.columns:
        pdf.project_name = project_data['ProjectName'].iloc[0]
    
    # Add report cover
    pdf.add_page()
    
    # Cover title
    pdf.set_font('Arial', 'B', 24)
    rgb = pdf.hex_to_rgb(ARCADIS_PRIMARY_COLOR)
    pdf.set_text_color(rgb[0], rgb[1], rgb[2])
    pdf.cell(0, 20, f"RiskLens Pro", 0, 1, 'C')
    
    pdf.set_font('Arial', 'B', 18)
    pdf.set_text_color(80, 80, 80)
    pdf.cell(0, 10, f"{report_type.title()} Report", 0, 1, 'C')
    
    pdf.ln(20)
    
    # Add Arcadis logo
    if os.path.exists('./assets/arcadis_logo.png'):
        pdf.image('./assets/arcadis_logo.png', x=75, y=pdf.get_y(), w=60)
    # Add business risk icon
    elif os.path.exists('./assets/business_risk_icon.png'):
        pdf.image('./assets/business_risk_icon.png', x=75, y=pdf.get_y(), w=60)
    
    pdf.ln(20)
    
    # Add date
    pdf.set_font('Arial', '', 12)
    pdf.cell(0, 10, f"Generated: {pd.Timestamp.now().strftime('%B %d, %Y')}", 0, 1, 'C')
    
    # Executive Summary
    pdf.add_page()
    pdf.chapter_title("Executive Summary")
    
    # Add executive summary text
    if 'summary_text' in risk_analysis:
        pdf.chapter_body(risk_analysis['summary_text'])
    else:
        pdf.chapter_body("This report provides an analysis of project risks based on historical data and "
                        "predictive modeling. It identifies potential high-risk projects and recommendations "
                        "for risk mitigation.")
    
    # Add key metrics
    pdf.ln(5)
    pdf.chapter_title("Key Risk Metrics", level=2)
    
    # Create a row of metrics
    metrics_per_row = 3
    metrics = [
        {"label": "Projects Analyzed", "value": len(project_data) if project_data is not None else 0},
        {"label": "High Risk Projects", "value": risk_analysis.get('high_risk_count', 0)},
        {"label": "Average Risk Score", "value": f"{risk_analysis.get('avg_risk_score', 0):.2f}"},
        {"label": "Risk Threshold", "value": f"{risk_analysis.get('risk_threshold', 0.5):.2f}"},
        {"label": "Model Accuracy", "value": f"{risk_analysis.get('model_accuracy', 0):.1f}%"},
        {"label": "Data Completeness", "value": f"{risk_analysis.get('data_completeness', 0):.1f}%"}
    ]
    
    for i, metric in enumerate(metrics):
        pdf.add_metric(metric["label"], metric["value"])
        
        # Add a line break every 'metrics_per_row' metrics
        if (i + 1) % metrics_per_row == 0:
            pdf.ln(30)
    
    # Add line break after metrics
    pdf.ln(30)
    
    # Add risk distribution chart
    pdf.chapter_title("Risk Distribution", level=2)
    
    if 'risk_distribution_fig' in visualizations:
        pdf.add_figure(visualizations['risk_distribution_fig'])
    
    # Add top risk factors
    pdf.chapter_title("Top Risk Factors", level=2)
    
    if 'risk_factors_fig' in visualizations:
        pdf.add_figure(visualizations['risk_factors_fig'])
    
    # Add recommendations section
    pdf.add_page()
    pdf.chapter_title("Risk Mitigation Recommendations", level=2)
    
    if 'recommendations' in risk_analysis:
        for i, rec in enumerate(risk_analysis['recommendations']):
            pdf.set_font('Arial', 'B', 11)
            pdf.set_text_color(70, 70, 70)
            pdf.cell(0, 8, f"{i+1}. {rec['title']}", 0, 1)
            
            pdf.set_font('Arial', '', 10)
            pdf.set_text_color(0, 0, 0)
            pdf.multi_cell(0, 5, rec['description'])
            pdf.ln(3)
    else:
        pdf.chapter_body("No specific recommendations available. Use the RiskLens Pro application to generate detailed recommendations.")
    
    # Add detailed sections based on report type
    if report_type == "detailed":
        # Add portfolio analysis
        pdf.add_page()
        pdf.chapter_title("Portfolio Analysis")
        
        # Add portfolio breakdown charts
        if 'portfolio_breakdown_fig' in visualizations:
            pdf.add_figure(visualizations['portfolio_breakdown_fig'])
        
        # Add top risky projects table
        if 'top_risky_projects' in risk_analysis and not risk_analysis['top_risky_projects'].empty:
            pdf.add_table(risk_analysis['top_risky_projects'], title="Top High-Risk Projects")
        
        # Add model analysis section
        pdf.add_page()
        pdf.chapter_title("Model Analysis")
        
        if 'model_metrics' in risk_analysis:
            pdf.chapter_body(f"The risk prediction model achieved {risk_analysis['model_metrics'].get('accuracy', 0):.1f}% accuracy "  
                            f"with {risk_analysis['model_metrics'].get('precision', 0):.2f} precision and "  
                            f"{risk_analysis['model_metrics'].get('recall', 0):.2f} recall.")
        
        # Add feature importance chart
        if 'feature_importance_fig' in visualizations:
            pdf.add_figure(visualizations['feature_importance_fig'])
        
        # Add Monte Carlo simulation results
        if 'monte_carlo_dist' in visualizations:
            pdf.add_page()
            pdf.chapter_title("Monte Carlo Simulation")
            pdf.chapter_body("Monte Carlo simulation was used to analyze the uncertainty in risk predictions "  
                            "and identify the most critical risk factors.")
            
            pdf.add_figure(visualizations['monte_carlo_dist'])
            
            if 'monte_carlo_metrics' in risk_analysis:
                pdf.ln(5)
                metrics_row = [
                    {"label": "Mean Risk", "value": f"{risk_analysis['monte_carlo_metrics'].get('mean', 0):.2f}"},
                    {"label": "Std Deviation", "value": f"{risk_analysis['monte_carlo_metrics'].get('std', 0):.2f}"},
                    {"label": "Conf. Interval", "value": f"[{risk_analysis['monte_carlo_metrics'].get('p10', 0):.2f} - {risk_analysis['monte_carlo_metrics'].get('p90', 0):.2f}]"}
                ]
                
                for metric in metrics_row:
                    pdf.add_metric(metric["label"], metric["value"])
    
    # Return PDF as bytes
    return pdf.output(dest='S').encode('latin1')


def create_pptx_report(project_data, risk_analysis, visualizations, report_type="executive"):
    """Create a PowerPoint presentation with project risk analysis
    
    Args:
        project_data: DataFrame with project data
        risk_analysis: Dictionary with risk analysis results
        visualizations: Dictionary with plotly figures
        report_type: Type of report to generate (executive, detailed, etc.)
        
    Returns:
        bytes: PowerPoint presentation as bytes object
    """
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "RiskLens Pro"
    subtitle.text = f"{report_type.title()} Report - {pd.Timestamp.now().strftime('%B %d, %Y')}"
    
    # Function to save a figure as an image and return a byte stream
    def fig_to_image(fig, width=1200, height=675):
        img_bytes = io.BytesIO()
        fig.write_image(img_bytes, format='png', width=width, height=height)
        img_bytes.seek(0)
        return img_bytes
    
    # Function to add a slide with title and content
    def add_content_slide(title_text, content_text=None, figure=None):
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        
        # Add content text if provided
        if content_text:
            content = slide.placeholders[1]
            content.text = content_text
        
        # Add figure if provided
        if figure:
            # If content text was provided, adjust figure position
            if content_text:
                # Save figure as image
                img_bytes = fig_to_image(figure)
                slide.shapes.add_picture(img_bytes, Inches(1), Inches(2.5), width=Inches(11))
            else:
                # Save figure as image
                img_bytes = fig_to_image(figure)
                slide.shapes.add_picture(img_bytes, Inches(1), Inches(1.5), width=Inches(11))
        
        return slide
    
    # Add executive summary slide
    summary_text = risk_analysis.get('summary_text', 
                                   "This report provides an analysis of project risks based on historical data and " 
                                   "predictive modeling. It identifies potential high-risk projects and recommendations " 
                                   "for risk mitigation.")
    
    add_content_slide("Executive Summary", summary_text)
    
    # Add risk distribution slide
    if 'risk_distribution_fig' in visualizations:
        add_content_slide("Risk Distribution", 
                         "Distribution of project risk probabilities across the portfolio.", 
                         visualizations['risk_distribution_fig'])
    
    # Add top risk factors slide
    if 'risk_factors_fig' in visualizations:
        add_content_slide("Top Risk Factors", 
                         "The most influential factors contributing to project risk.", 
                         visualizations['risk_factors_fig'])
    
    # Add recommendations slide
    if 'recommendations' in risk_analysis and risk_analysis['recommendations']:
        rec_slide = add_content_slide("Risk Mitigation Recommendations")
        content = rec_slide.placeholders[1]
        
        rec_text = ""
        for i, rec in enumerate(risk_analysis['recommendations'][:5]):  # Limit to top 5 recommendations
            rec_text += f"{i+1}. {rec['title']}\n"
            rec_text += f"   {rec['description']}\n\n"
        
        content.text = rec_text
    
    # Add detailed slides based on report type
    if report_type == "detailed":
        # Add portfolio analysis slide
        if 'portfolio_breakdown_fig' in visualizations:
            add_content_slide("Portfolio Analysis", 
                             "Breakdown of risk across different project categories.", 
                             visualizations['portfolio_breakdown_fig'])
        
        # Add model analysis slide
        if 'feature_importance_fig' in visualizations:
            # Create model metrics text
            if 'model_metrics' in risk_analysis:
                metrics_text = f"Model Performance:\n\n" \
                               f"Accuracy: {risk_analysis['model_metrics'].get('accuracy', 0):.1f}%\n" \
                               f"Precision: {risk_analysis['model_metrics'].get('precision', 0):.2f}\n" \
                               f"Recall: {risk_analysis['model_metrics'].get('recall', 0):.2f}\n" \
                               f"\nThe chart shows the relative importance of different features in predicting project risk."
            else:
                metrics_text = "The chart shows the relative importance of different features in predicting project risk."
            
            add_content_slide("Model Analysis", metrics_text, visualizations['feature_importance_fig'])
        
        # Add Monte Carlo simulation slide
        if 'monte_carlo_dist' in visualizations:
            mc_text = "Monte Carlo simulation results showing the distribution of risk probabilities " \
                     "across thousands of simulated scenarios.\n\n"
            
            if 'monte_carlo_metrics' in risk_analysis:
                mc_text += f"Mean Risk: {risk_analysis['monte_carlo_metrics'].get('mean', 0):.2f}\n" \
                          f"Std Deviation: {risk_analysis['monte_carlo_metrics'].get('std', 0):.2f}\n" \
                          f"80% Confidence Interval: [{risk_analysis['monte_carlo_metrics'].get('p10', 0):.2f} - {risk_analysis['monte_carlo_metrics'].get('p90', 0):.2f}]"
            
            add_content_slide("Monte Carlo Simulation", mc_text, visualizations['monte_carlo_dist'])
    
    # Save to a byte stream
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes.getvalue()


def get_download_link(binary_data, filename, text):
    """Generate a download link for binary data"""
    b64 = base64.b64encode(binary_data).decode()
    href = f'<a href="data:application/octet-stream;base64,{b64}" download="{filename}">{text}</a>'
    return href


def create_report_download_section(project_data, risk_analysis, visualizations):
    """Create a section with report download options"""
    st.markdown("### Download Reports")
    st.markdown("Generate and download comprehensive reports with all visualizations and analysis.")
    
    # Report type selection
    report_type = st.radio(
        "Report Type:",
        options=["Executive Summary", "Detailed Report"],
        horizontal=True,
        help="Executive Summary provides a high-level overview, while Detailed Report includes all analyses."
    )
    
    report_type_value = "executive" if report_type == "Executive Summary" else "detailed"
    
    # Create two columns for PDF and PPTX downloads
    col1, col2 = st.columns(2)
    
    with col1:
        pdf_button = st.button("Generate PDF Report", key="pdf_button", use_container_width=True)
        
        if pdf_button:
            with st.spinner("Generating PDF report..."):
                pdf_bytes = create_pdf_report(project_data, risk_analysis, visualizations, report_type_value)
                st.markdown(
                    get_download_link(pdf_bytes, f"risklens_pro_{report_type_value}_report.pdf", "Download PDF Report"),
                    unsafe_allow_html=True
                )
    
    with col2:
        pptx_button = st.button("Generate PowerPoint", key="pptx_button", use_container_width=True)
        
        if pptx_button:
            with st.spinner("Generating PowerPoint presentation..."):
                pptx_bytes = create_pptx_report(project_data, risk_analysis, visualizations, report_type_value)
                st.markdown(
                    get_download_link(pptx_bytes, f"risklens_pro_{report_type_value}_presentation.pptx", "Download PowerPoint"),
                    unsafe_allow_html=True
                )
